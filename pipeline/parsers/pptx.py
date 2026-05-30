"""
pipeline/parsers/pptx.py

Parse .pptx files. One "page" per slide.
Slide body text and speaker notes are combined; notes are tagged with a
[Speaker Notes] marker so downstream agents can weight them differently.
"""

from pptx import Presentation


def extract(filepath: str) -> list[dict]:
    prs = Presentation(filepath)
    pages = []

    for i, slide in enumerate(prs.slides, start=1):
        body_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    body_parts.append(text)

        notes_text = ""
        if slide.has_notes_slide:
            nf = slide.notes_slide.notes_text_frame
            if nf and nf.text.strip():
                notes_text = nf.text.strip()

        combined = "\n".join(body_parts)
        if notes_text:
            combined = (
                combined + "\n\n[Speaker Notes]\n" + notes_text).strip()

        if combined:
            pages.append({
                "page": i,
                "text": combined,
                "language": None,
            })

    return pages
